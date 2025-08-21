
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml import parse_xml
from pptx.util import Cm
import io, requests, os, random, textwrap, datetime

# Optional Unsplash integration
UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY")

def rgb_from_hex(hx: str) -> RGBColor:
    hx = hx.strip().replace("#", "")
    if len(hx) != 6:
        hx = "FFFFFF"
    r = int(hx[0:2], 16); g = int(hx[2:4], 16); b = int(hx[4:6], 16)
    return RGBColor(r, g, b)

def add_textbox(slide, left_cm, top_cm, width_cm, height_cm, text, font_name="Calibri", font_size=24, color="#1A1A1A", bold=False):
    left, top, width, height = Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    if bold:
        run.font.bold = True
    run.font.size = Pt(font_size)
    run.font.name = font_name
    run.font.color.rgb = rgb_from_hex(color)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return shape

def add_bullets(slide, left_cm, top_cm, width_cm, height_cm, bullets, font_name="Calibri", font_size=18, color="#1A1A1A"):
    left, top, width, height = Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = rgb_from_hex(color)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return shape

def fetch_image_bytes(keyword: str) -> bytes:
    try:
        if UNSPLASH_ACCESS_KEY:
            # Basic Unsplash search
            import urllib.parse, random, requests
            q = urllib.parse.quote(keyword or "creative")
            url = f"https://api.unsplash.com/photos/random?query={q}&orientation=landscape&content_filter=high"
            headers = {"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"}
            r = requests.get(url, headers=headers, timeout=15)
            r.raise_for_status()
            data = r.json()
            img_url = data.get("urls", {}).get("regular") or data.get("urls", {}).get("full")
        else:
            # Fallback to picsum (no key)
            img_url = "https://picsum.photos/1600/900"
        img = requests.get(img_url, timeout=20)
        img.raise_for_status()
        return img.content
    except Exception:
        # solid color fallback
        from PIL import Image
        import io as _io
        im = Image.new("RGB", (1600, 900), (230, 230, 230))
        bio = _io.BytesIO()
        im.save(bio, format="PNG")
        return bio.getvalue()

def layout_title_only(prs, slide, spec):
    # Title big in center-ish
    pass

def apply_background(slide, color_hex: str):
    if not color_hex: 
        return
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_from_hex(color_hex)

def add_image(slide, img_bytes, left_cm, top_cm, width_cm=None, height_cm=None):
    from pptx.util import Cm
    pic = slide.shapes.add_picture(io.BytesIO(img_bytes), Cm(left_cm), Cm(top_cm), Cm(width_cm) if width_cm else None, Cm(height_cm) if height_cm else None)
    return pic

def add_speaker_notes(slide, notes_text: str):
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = notes_text

def pick_layout_name(style_layout: str):
    # Map style layout to handlers
    mapping = {
        "title_only": "_title_only",
        "title_content": "_title_content",
        "image_left_text_right": "_img_left_text_right",
        "text_left_image_right": "_text_left_img_right",
        "quote": "_quote",
        "two_columns": "_two_columns",
        "comparison": "_comparison",
        "timeline": "_timeline",
        "big_number": "_big_number",
        "chart_focus": "_chart_focus"
    }
    return mapping.get(style_layout, "_title_content")

def render_slide(prs, slide, slide_spec, seed=None):
    rnd = random.Random(seed)
    title = slide_spec.get("title", "")
    bullets = slide_spec.get("bullets", [])
    notes = slide_spec.get("speaker_notes", "")
    style = slide_spec.get("style", {})
    bg = style.get("background_color", "#FFFFFF")
    title_color = style.get("title_color", "#111111")
    text_color = style.get("text_color", "#222222")
    font = style.get("font_family", "Calibri")
    layout_name = pick_layout_name(style.get("layout", "title_content"))
    has_image = style.get("has_image", False)
    image_keyword = style.get("image_keyword", "")

    apply_background(slide, bg)

    # Title position varies slightly
    title_top = rnd.uniform(0.8, 2.0)
    add_textbox(slide, left_cm=1.5, top_cm=title_top, width_cm=22, height_cm=2.5,
                text=title, font_name=font, font_size=34, color=title_color, bold=True)

    # Layout rendering
    if layout_name == "_title_content":
        add_bullets(slide, left_cm=2, top_cm=title_top+2.3, width_cm=22, height_cm=11,
                    bullets=bullets, font_name=font, font_size=20, color=text_color)
        if has_image:
            img = fetch_image_bytes(image_keyword)
            add_image(slide, img, left_cm=18, top_cm=title_top+2.5, width_cm=7.0, height_cm=None)

    elif layout_name == "_img_left_text_right":
        if has_image:
            img = fetch_image_bytes(image_keyword)
            add_image(slide, img, left_cm=1.2, top_cm=title_top+2.5, width_cm=11.5, height_cm=None)
        add_bullets(slide, left_cm=13.2, top_cm=title_top+2.5, width_cm=11, height_cm=11,
                    bullets=bullets, font_name=font, font_size=20, color=text_color)

    elif layout_name == "_text_left_img_right":
        add_bullets(slide, left_cm=1.2, top_cm=title_top+2.5, width_cm=11.5, height_cm=11,
                    bullets=bullets, font_name=font, font_size=20, color=text_color)
        if has_image:
            img = fetch_image_bytes(image_keyword)
            add_image(slide, img, left_cm=13.2, top_cm=title_top+2.5, width_cm=11, height_cm=None)

    elif layout_name == "_quote":
        quote = "\\n".join(bullets) if bullets else ""
        add_textbox(slide, left_cm=3, top_cm=title_top+2.6, width_cm=18, height_cm=8,
                    text=f"“{quote}”", font_name=font, font_size=30, color=text_color, bold=False)

    elif layout_name == "_two_columns":
        mid = 12.5
        left_bul = bullets[:len(bullets)//2 or 1]
        right_bul = bullets[len(bullets)//2 or 1:]
        add_bullets(slide, left_cm=1.5, top_cm=title_top+2.5, width_cm=10.5, height_cm=11,
                    bullets=left_bul, font_name=font, font_size=19, color=text_color)
        add_bullets(slide, left_cm=mid, top_cm=title_top+2.5, width_cm=10.5, height_cm=11,
                    bullets=right_bul, font_name=font, font_size=19, color=text_color)

    elif layout_name == "_comparison":
        labels = bullets if bullets else ["Ưu điểm", "Nhược điểm"]
        add_textbox(slide, 2, title_top+2.4, 10.5, 1, "So sánh", font, 18, text_color, True)
        add_bullets(slide, 2, title_top+3.4, 10.5, 10, labels[:len(labels)//2 or 1], font, 18, text_color)
        add_bullets(slide, 13.2, title_top+3.4, 10.5, 10, labels[len(labels)//2 or 1:], font, 18, text_color)

    elif layout_name == "_timeline":
        # simple vertical timeline using bullets
        y = title_top+2.3
        for i, b in enumerate(bullets[:6]):
            add_textbox(slide, 2, y + i*1.8, 1, 1, f"{i+1}", font, 18, text_color, True)
            add_textbox(slide, 3.2, y + i*1.8, 20.5, 1.6, b, font, 18, text_color, False)

    elif layout_name == "_big_number":
        num = next((x for x in bullets if any(c.isdigit() for c in x)), "100%")
        add_textbox(slide, 2, title_top+4.0, 21, 6, num, font, 72, text_color, True)
        remain = [x for x in bullets if x != num]
        if remain:
            add_bullets(slide, 2, title_top+8.2, 21, 4, remain, font, 20, text_color)

    elif layout_name == "_chart_focus":
        # Placeholder: indicate where chart should be
        add_textbox(slide, 2, title_top+2.5, 10, 1, "Biểu đồ (gợi ý):", font, 18, text_color, True)
        add_bullets(slide, 2, title_top+3.5, 10, 9, ["Cột", "Đường", "Bánh"], font, 18, text_color)
        add_textbox(slide, 13.2, title_top+2.5, 10.5, 10, "\\n".join(bullets), font, 18, text_color)

    add_speaker_notes(slide, notes)

def create_presentation_from_spec(spec: dict, seed=None):
    rnd = random.Random(seed if seed is not None else random.randint(0, 10_000_000))
    prs = Presentation()
    prs.slide_width = Cm(33.867)  # 16:9 default in centimeters
    prs.slide_height = Cm(19.05)

    title = spec.get("title", "AI Generated Slides")
    slides = spec.get("slides", [])

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank, we do custom
    apply_background(slide, "#FFFFFF")
    add_textbox(slide, 2, 5.0, 26, 3.2, title, "Georgia", 44, "#111111", True)
    add_textbox(slide, 2, 9.0, 24, 2, "Tạo bởi AI Slide Designer", "Calibri", 20, "#444444")

    # Deck slides
    for i, s in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        render_slide(prs, slide, s, seed=rnd.randint(0, 10_000_000))

    # Save to bytes
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    fname = f"slides.{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    return bio.read(), fname
